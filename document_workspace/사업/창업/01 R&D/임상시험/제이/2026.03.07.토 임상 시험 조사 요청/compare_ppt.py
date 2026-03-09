import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

def extract_pptx(filepath, label):
    prs = Presentation(filepath)
    print(f'\n{"="*60}')
    print(f'{label}')
    print(f'슬라이드 수: {len(prs.slides)}')
    print(f'슬라이드 크기: {prs.slide_width/914400:.1f}x{prs.slide_height/914400:.1f} inches')
    print(f'{"="*60}')
    for i, slide in enumerate(prs.slides):
        print(f'\n--- 슬라이드 {i+1} ---')
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if len(text) > 600:
                        text = text[:600] + '...(truncated)'
                    print(f'  [TEXT] {text}')
            if shape.has_table:
                table = shape.table
                print(f'  [TABLE {table.rows.__len__()}x{len(table.columns)}]')
                for r in range(min(table.rows.__len__(), 20)):
                    row_text = []
                    for c in range(len(table.columns)):
                        cell_text = table.cell(r, c).text.strip().replace('\n', ' / ')
                        row_text.append(cell_text)
                    print(f'    R{r}: {" | ".join(row_text)}')

extract_pptx('비임상_임상_파이프라인_20260308_1555_코워크.pptx', '코워크 버전')
print('\n\n')
extract_pptx('치료제_비임상_임상_파이프라인_20260308_1550_클로드.pptx', '클로드 버전')
