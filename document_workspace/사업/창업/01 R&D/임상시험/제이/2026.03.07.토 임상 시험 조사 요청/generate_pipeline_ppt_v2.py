"""
치료제 비임상/임상 파이프라인 PPT 생성 스크립트 v2
- 코워크 색상, 표준 16:9 (13.333 x 7.5)
- 각 페이지: 간트차트 + 상세 테이블 (기간, 비용, 업체)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

S = 4 / 3


def I(val):
    return Inches(val * S)


# ── 코워크 색상 팔레트 ──
CW_DARK_BG = RGBColor(0x0B, 0x20, 0x27)
CW_TEAL_BRIGHT = RGBColor(0x02, 0xC3, 0x9A)
CW_TEAL = RGBColor(0x00, 0xA8, 0x96)
CW_TEAL_DARK = RGBColor(0x02, 0x80, 0x90)
CW_TEAL_DARKER = RGBColor(0x02, 0x6B, 0x78)
CW_RED = RGBColor(0xE8, 0x44, 0x3A)
CW_AMBER = RGBColor(0xF4, 0xA2, 0x61)
CW_CONTENT_BG = RGBColor(0xF0, 0xF9, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
TABLE_ROW_EVEN = RGBColor(0xE8, 0xF4, 0xF2)
TABLE_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=10, bold=False,
                 color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, size=10, bold=False, color=DARK_GRAY,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(1)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = alignment
    p.space_before = space_before
    return p


def make_table(slide, left, top, width, rows, cols, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.1))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    return table_shape, table


def style_cell(cell, text, size=7, bold=False, color=DARK_GRAY, fill=None,
               alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(3)
    cell.margin_bottom = Pt(3)


# ═══════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, CW_DARK_BG)

add_rect(s1, I(0), I(0), I(10), I(0.06), CW_TEAL_BRIGHT)
add_rect(s1, I(0.7), I(2.061), I(0.06), I(1.503), CW_TEAL)

tf_title = add_text_box(s1, I(1.1), I(2.013), I(7), I(1.6),
                         "치료제 비임상 / 임상", size=32 * S, bold=True, color=WHITE)
add_paragraph(tf_title, "파이프라인 정리", size=32 * S, bold=True, color=WHITE,
              space_before=Pt(6))

add_rect(s1, I(0), I(5.325), I(10), I(0.3), CW_TEAL_DARK)
add_text_box(s1, I(0.5), I(5.325), I(9), I(0.3),
             "근거: 업체 미팅 기록 (Dt&CRO, 에피바이오텍, P&K, Bio FD&C) + 인터넷 크로스체크",
             size=9 * S, color=RGBColor(0xCC, 0xDD, 0xDD), alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# 슬라이드 2: 비임상 (Preclinical) — 간트차트 + 테이블
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, CW_CONTENT_BG)

# 상단 헤더바
add_rect(s2, I(0), I(0), I(10), I(0.65), CW_TEAL_DARK)
add_text_box(s2, I(0.5), I(0.12), I(7), I(0.45),
             "비임상 (Preclinical) 파이프라인",
             size=20 * S, bold=True, color=WHITE)
add_text_box(s2, I(7.5), I(0.12), I(2), I(0.45),
             "예상 ~3년", size=11 * S, bold=True, color=RGBColor(0xCC, 0xEE, 0xE8),
             alignment=PP_ALIGN.RIGHT)

# ── 간트차트 (균등 년차) ──
GANTT_LABEL_L = 0.4
GANTT_CHART_L = 1.7
GANTT_CHART_W = 7.9   # 3년 = 7.9인치 → 균등
GANTT_TOP = 0.85
GANTT_BAR_H = I(0.15)
GANTT_BAR_GAP = I(0.19)

year_w = GANTT_CHART_W / 3.0
NUM_BARS = 7

# 흰색 컨테이너
container_h = GANTT_BAR_GAP * NUM_BARS + I(0.35)
add_rect(s2, I(GANTT_LABEL_L - 0.05), I(GANTT_TOP - 0.05),
         I(GANTT_CHART_W + (GANTT_CHART_L - GANTT_LABEL_L) + 0.15),
         container_h, WHITE)

# 세로 그리드 + 년차 라벨
for yr in range(4):
    x = GANTT_CHART_L + year_w * yr
    if yr > 0:
        add_rect(s2, I(x), I(GANTT_TOP), Inches(1),
                 GANTT_BAR_GAP * NUM_BARS + I(0.05), LIGHT_GRAY)
    label = f"Y{yr}" if yr > 0 else "0"
    add_text_box(s2, I(x - 0.15), I(GANTT_TOP) + GANTT_BAR_GAP * NUM_BARS + I(0.08),
                 I(0.5), I(0.15),
                 label, size=7 * S, bold=True, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)

# 간트 바 (테이블 7단계와 일치)
CW_TEAL_LIGHT = RGBColor(0x48, 0xC9, 0xB0)  # 기전 연구용
CW_BLUE = RGBColor(0x3A, 0x7C, 0xA5)        # PK 연구용

gantt_items = [
    ("① 분석법 개발", 0.0,  1.0,  CW_TEAL,       "분석법 개발(6개월~1년)"),
    ("② 효력 시험",   0.5,  0.5,  CW_TEAL_BRIGHT, "효력 시험(2 ~ 8주)"),
    ("③ 기전 연구",   0.5,  0.5,  CW_TEAL_LIGHT,  "기전 연구(효력 시험 병행)"),
    ("④ GLP 독성",   1.0,  1.5,  CW_RED,         "GLP 독성(1 ~ 1.5년)"),
    ("⑤ PK 연구",    1.0,  1.5,  CW_BLUE,        "PK 연구(독성 시험 병행)"),
    ("⑥ CMC/제형",   0.5,  2.0,  CW_TEAL_DARK,   "CMC/제형(0.5 ~ 2.5년)"),
    ("⑦ IND 준비",   2.5,  0.5,  CW_AMBER,       "IND 준비"),
]

for i, (label, start_yr, dur_yr, clr, bar_text) in enumerate(gantt_items):
    y = I(GANTT_TOP + 0.05) + GANTT_BAR_GAP * i
    add_text_box(s2, I(GANTT_LABEL_L), y, I(1.2), GANTT_BAR_H,
                 label, size=6.5 * S, bold=True, color=DARK_GRAY,
                 alignment=PP_ALIGN.RIGHT)
    bar_x = GANTT_CHART_L + year_w * start_yr
    bar_w = year_w * dur_yr
    add_rect(s2, I(bar_x), y, I(bar_w), GANTT_BAR_H, clr)
    add_text_box(s2, I(bar_x + 0.05), y, I(bar_w - 0.1), GANTT_BAR_H,
                 bar_text, size=5.5 * S, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

# ── 단계별 상세 테이블 (전체 너비 사용) ──
tbl_top = 2.55
tbl_headers = ["단계", "주요 내용", "기간", "예상 비용", "협력 업체"]
tbl_data = [
    ["① 분석법 개발", "HPLC/Mass 분석법 + 생체시료 분석법 확립", "6개월~1년", "별도 협의", "Dt&CRO"],
    ["② 효력 시험", "양모 모델 + 남성형 탈모 질환 모델(DHT) / 미녹시딜 양성대조", "2~8주", "1,500~3,000만", "에피바이오텍"],
    ["③ 기전 연구", "바이오마커 qPCR, 조직 분석 / 모낭 직경/밀도/무게 측정", "효력시험 병행", "효력+기전 / ~3,000만", "에피바이오텍"],
    ["④ GLP 독성", "단회→DRF(4주)→본시험(13주) / + 유전독성", "1~1.5년", "8억~12억", "Dt&CRO"],
    ["⑤ PK 연구", "ADME, 경피 특성 / 투여 부위 조직 분석", "독성시험 병행", "독성 패키지 포함", "Dt&CRO"],
    ["⑥ CMC", "원료 분석, 품질관리, 안정성, 제형 스케일업", "병행 진행", "1~2억 (추정)", "미정"],
    ["⑦ IND 준비", "비임상 데이터 통합 + 임상시험계획서 / 식약처 사전상담/제출", "3~6개월", "RA/CRO 비용", "미정"],
]

tbl_col_w = [I(1.1), I(4.0), I(1.0), I(1.4), I(1.2)]
tbl_shape, tbl = make_table(s2, I(0.35), I(tbl_top), I(9.3),
                             len(tbl_data)+1, len(tbl_headers), tbl_col_w)

for j, h in enumerate(tbl_headers):
    style_cell(tbl.cell(0, j), h, size=8 * S, bold=True, color=WHITE,
               fill=CW_TEAL_DARK, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(tbl_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    for j, val in enumerate(row):
        c = CW_RED if val == "미정" else DARK_GRAY
        style_cell(tbl.cell(i+1, j), val, size=7 * S, color=c, fill=bg)

# ── 비임상 협력 업체 ──
pre_partner_top = 4.55
add_rect(s2, I(0.35), I(pre_partner_top), I(9.3), I(0.22), CW_TEAL_DARK)
add_text_box(s2, I(0.45), I(pre_partner_top), I(4), I(0.22),
             "비임상 협력 업체",
             size=9 * S, bold=True, color=WHITE)

pre_partner_headers = ["업체", "역할", "미팅일", "상태", "핵심 메모"]
pre_partner_data = [
    ["에피바이오텍", "비임상 효력 (탈모 CRO)", "2025.12", "완료", "견적서 수령 예정, 남성형 질환모델 보유"],
    ["Dt&CRO", "GLP 독성 + 분석법", "2025.11", "완료", "4주 ~8억, 13주 ~10-12억"],
    ["Bio FD&C", "OEM/ODM + INCI", "2025.10", "완료", "INCI 원료등록 신청 완료 (2026.02)"],
    ["CMC 전문업체", "원료 품질관리", "-", "미탐색", ""],
    ["RA 전문가", "인허가 전략 (IND)", "-", "미탐색", ""],
]

pre_p_col_w = [I(1.1), I(1.8), I(0.7), I(0.8), I(4.9)]
pre_p_shape, pre_p_tbl = make_table(s2, I(0.35), I(pre_partner_top + 0.22), I(9.3),
                                     len(pre_partner_data)+1, 5, pre_p_col_w)

for j, h in enumerate(pre_partner_headers):
    style_cell(pre_p_tbl.cell(0, j), h, size=7 * S, bold=True, color=WHITE,
               fill=CW_TEAL_DARKER, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(pre_partner_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    for j, val in enumerate(row):
        if val == "미탐색":
            c = CW_RED
        elif val == "완료":
            c = CW_TEAL
        else:
            c = DARK_GRAY
        style_cell(pre_p_tbl.cell(i+1, j), val, size=6.5 * S, color=c, fill=bg,
                   alignment=PP_ALIGN.CENTER if j in [2, 3] else PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# 슬라이드 3: 임상 (Clinical) — 간트차트 + 테이블
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, CW_CONTENT_BG)

# 상단 헤더바
add_rect(s3, I(0), I(0), I(10), I(0.65), CW_TEAL_DARK)
add_text_box(s3, I(0.5), I(0.12), I(7), I(0.45),
             "임상 (Clinical) 파이프라인",
             size=20 * S, bold=True, color=WHITE)
add_text_box(s3, I(7.5), I(0.12), I(2), I(0.45),
             "예상 ~7년", size=11 * S, bold=True, color=RGBColor(0xCC, 0xEE, 0xE8),
             alignment=PP_ALIGN.RIGHT)

# ── 임상 간트차트 ──
CLIN_LABEL_L = 0.4
CLIN_CHART_L = 1.7
CLIN_CHART_W = 7.9
CLIN_YEARS = 9
CLIN_TOP = 0.85
CLIN_BAR_H = I(0.22)
CLIN_BAR_GAP = I(0.32)

clin_year_w = CLIN_CHART_W / CLIN_YEARS

# 흰색 컨테이너
clin_container_h = CLIN_BAR_GAP * 4 + I(0.35)
add_rect(s3, I(CLIN_LABEL_L - 0.05), I(CLIN_TOP - 0.05),
         I(CLIN_CHART_W + (CLIN_CHART_L - CLIN_LABEL_L) + 0.15),
         clin_container_h, WHITE)

# 년차 라벨
for yr in range(CLIN_YEARS + 1):
    x = CLIN_CHART_L + clin_year_w * yr
    if 0 < yr < CLIN_YEARS:
        add_rect(s3, I(x), I(CLIN_TOP), Inches(1),
                 CLIN_BAR_GAP * 4 + I(0.05), LIGHT_GRAY)
    label = f"{yr}yr" if yr > 0 else "0"
    add_text_box(s3, I(x - 0.2), I(CLIN_TOP) + CLIN_BAR_GAP * 4 + I(0.08),
                 I(0.5), I(0.15),
                 label, size=6 * S, bold=True, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)

# 임상 간트 바
clin_gantt = [
    ("Phase 1",  0.0, 1.5, CW_TEAL,      "안전성+PK / 20~80명"),
    ("Phase 2",  1.5, 2.5, CW_TEAL_DARK, "AGA 100~250명 / 용량탐색"),
    ("Phase 3",  4.0, 3.0, CW_RED,       "대규모 300~3,000명"),
    ("NDA 허가", 7.0, 2.0, CW_AMBER,     "품목허가 심사"),
]

for i, (label, start_yr, dur_yr, clr, bar_text) in enumerate(clin_gantt):
    y = I(CLIN_TOP + 0.05) + CLIN_BAR_GAP * i
    add_text_box(s3, I(CLIN_LABEL_L), y, I(1.2), CLIN_BAR_H,
                 label, size=8 * S, bold=True, color=DARK_GRAY,
                 alignment=PP_ALIGN.RIGHT)
    bar_x = CLIN_CHART_L + clin_year_w * start_yr
    bar_w = clin_year_w * dur_yr
    add_rect(s3, I(bar_x), y, I(bar_w), CLIN_BAR_H, clr)
    add_text_box(s3, I(bar_x + 0.05), y, I(bar_w - 0.1), CLIN_BAR_H,
                 bar_text, size=6 * S, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

# ── 임상 단계별 상세 테이블 (전체 너비) ──
clin_tbl_top = 2.65
clin_headers = ["단계", "주요 내용", "기간", "예상 비용", "협력 업체/비고"]
clin_data = [
    ["Phase 1", "안전성 + PK / 건강인·환자 20~80명 / 외용제: Phase 1/2 통합 가능", "1~1.5년", "수억~10억+", "OPIS"],
    ["Phase 2", "AGA 환자 유효성·안전성 / 100~250명 / 용량 탐색", "2~2.5년", "수십억 원", "OPIS"],
    ["Phase 3", "대규모 확증 시험 / 300~3,000명 / 미녹시딜·경구 Fina 대조", "2.5~3년", "수백억 원", "OPIS / 대웅제약"],
    ["NDA 허가", "품목허가 신청 → 식약처/FDA 심사", "1~2년", "RA 비용", "미정"],
]

clin_col_w = [I(0.9), I(4.3), I(1.0), I(1.3), I(1.8)]
clin_tbl_shape, clin_tbl = make_table(s3, I(0.35), I(clin_tbl_top), I(9.3),
                                       len(clin_data)+1, len(clin_headers), clin_col_w)

for j, h in enumerate(clin_headers):
    style_cell(clin_tbl.cell(0, j), h, size=8 * S, bold=True, color=WHITE,
               fill=CW_TEAL_DARK, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(clin_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    for j, val in enumerate(row):
        style_cell(clin_tbl.cell(i+1, j), val, size=7 * S, color=DARK_GRAY, fill=bg)

# ── 임상 협력 업체 ──
partner_top = 3.85
add_rect(s3, I(0.35), I(partner_top), I(9.3), I(0.22), CW_TEAL_DARK)
add_text_box(s3, I(0.45), I(partner_top), I(4), I(0.22),
             "임상 협력 업체",
             size=9 * S, bold=True, color=WHITE)

partner_headers = ["업체", "역할", "미팅일", "상태", "핵심 메모"]
partner_data = [
    ["P&K", "인체적용시험", "2025.12", "완료", "기능성 화장품 인증, 피부흡수도 평가"],
    ["OPIS (이탈리아)", "글로벌 CRO (I~IV)", "2026.02", "완료", "30개국 네트워크, EU 임상 협력"],
    ["대웅제약", "전략적 파트너", "-", "진행 중", "이노베어 5기, 기술이전 우선협상"],
]

p_col_w = [I(1.1), I(1.6), I(0.7), I(0.8), I(5.1)]
p_shape, p_tbl = make_table(s3, I(0.35), I(partner_top + 0.22), I(9.3),
                             len(partner_data)+1, 5, p_col_w)

for j, h in enumerate(partner_headers):
    style_cell(p_tbl.cell(0, j), h, size=7 * S, bold=True, color=WHITE,
               fill=CW_TEAL_DARKER, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(partner_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    for j, val in enumerate(row):
        if val == "미탐색":
            c = CW_RED
        elif val in ("완료", "진행 중"):
            c = CW_TEAL
        else:
            c = DARK_GRAY
        style_cell(p_tbl.cell(i+1, j), val, size=6.5 * S, color=c, fill=bg,
                   alignment=PP_ALIGN.CENTER if j in [2, 3] else PP_ALIGN.LEFT)


# ── 저장 ──
output_dir = os.path.dirname(os.path.abspath(__file__))
timestamp = datetime.now().strftime("%Y%m%d_%H%M")
output_path = os.path.join(output_dir, f"치료제_비임상_임상_파이프라인_{timestamp}_클로드.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
