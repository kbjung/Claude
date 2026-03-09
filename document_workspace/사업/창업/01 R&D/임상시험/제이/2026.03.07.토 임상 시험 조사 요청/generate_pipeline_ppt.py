"""
치료제 비임상/임상 파이프라인 PPT 생성 스크립트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
MID_BLUE = RGBColor(0x3A, 0x5B, 0x9F)
LIGHT_BLUE = RGBColor(0x5B, 0x8D, 0xD9)
ACCENT_ORANGE = RGBColor(0xE8, 0x7C, 0x2A)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
TABLE_HEADER_BG = RGBColor(0x2C, 0x3E, 0x6B)
TABLE_ROW_EVEN = RGBColor(0xEA, 0xEF, 0xF7)
TABLE_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)
WARN_BG = RGBColor(0xFF, 0xF3, 0xE0)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_text_box(slide, left, top, width, height, text, size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, space_before=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = alignment
    p.space_before = space_before
    return p


def make_table(slide, left, top, width, rows, cols, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.1))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    # Auto-adjust height later
    return table_shape, table


def style_cell(cell, text, size=8, bold=False, color=DARK_GRAY, fill=None, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    # margins
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)


# ═══════════════════════════════════════════════════════════════
# 페이지 1: 비임상 (Preclinical) 파이프라인
# ═══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, WHITE)

# ── 상단 바 ──
header_bar = add_rect(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
set_text(header_bar, "비임상 (Preclinical) 파이프라인", size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
header_bar.text_frame.paragraphs[0].space_before = Pt(8)

add_text_box(slide1, Inches(0.5), Inches(0.6), Inches(5), Inches(0.3),
             "펩타이드 기반 국소 Finasteride 탈모 치료제", size=11, color=RGBColor(0xBB,0xCC,0xEE))

# ── 타임라인 섹션 ──
add_text_box(slide1, Inches(0.4), Inches(1.0), Inches(3), Inches(0.3),
             "▎ 비임상 타임라인 (~3년)", size=13, bold=True, color=NAVY)

# 타임라인 bars
timeline_items = [
    ("분석법 개발",     0.0, 3.0, MID_BLUE),
    ("효력 시험 (동물)", 1.5, 3.0, ACCENT_GREEN),
    ("기전 연구",       2.0, 3.0, RGBColor(0x2E, 0xCC, 0x71)),
    ("GLP 독성 시험",   3.5, 5.5, ACCENT_ORANGE),
    ("PK 연구",        4.0, 4.5, RGBColor(0xF3, 0x9C, 0x12)),
    ("CMC/제형",       1.5, 6.5, LIGHT_BLUE),
    ("IND 준비",       7.5, 2.0, ACCENT_RED),
]

bar_left_base = Inches(2.5)
bar_top_base = Inches(1.4)
bar_height = Inches(0.25)
year_width = Inches(1.3)  # 1 year = 1.3 inches
bar_gap = Inches(0.30)

for i, (label, start, duration, color) in enumerate(timeline_items):
    y = bar_top_base + bar_gap * i
    # label
    add_text_box(slide1, Inches(0.4), y - Inches(0.02), Inches(2.0), Inches(0.25),
                 label, size=8, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT)
    # bar
    bar = add_rect(slide1, bar_left_base + Emu(int(start * year_width)), y,
                   Emu(int(duration * year_width)), bar_height, color)
    bar.line.fill.background()

# Year markers
for yr in range(4):
    x = bar_left_base + Emu(int(yr * 3.0 * year_width))
    add_text_box(slide1, x - Inches(0.2), bar_top_base + bar_gap * 7 + Inches(0.05),
                 Inches(0.8), Inches(0.2),
                 f"Y{yr+1}", size=8, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── 단계별 상세 테이블 ──
add_text_box(slide1, Inches(0.4), Inches(3.65), Inches(3), Inches(0.3),
             "▎ 단계별 상세", size=13, bold=True, color=NAVY)

headers = ["단계", "주요 내용", "기간", "예상 비용", "담당 업체", "상태"]
data_rows = [
    ["① 분석법 개발", "물질 분석법(HPLC/Mass)\n+ 생체시료 분석법 확립", "6개월~1년", "별도 협의", "Dt&CRO", "✅ 미팅완료"],
    ["② 효력 시험", "양모 모델 + 남성형 탈모\n질환 모델 (DHT 투여)", "2~8주", "1,500~3,000만", "에피바이오텍", "✅ 미팅완료"],
    ["③ 기전 연구", "바이오마커 qPCR\n조직/모낭 분석", "효력과 병행", "효력 합산\n~3,000만", "에피바이오텍", "✅ 미팅완료"],
    ["④ GLP 독성", "단회 → DRF(4주)\n→ 반복(13주) + 유전독성", "1~1.5년", "8억(4주)\n10~12억(13주)", "Dt&CRO", "✅ 미팅완료"],
    ["⑤ PK 연구", "ADME. 경피 특성상\n투여부위 조직 분석 병행", "독성과 병행", "독성 패키지\n포함", "Dt&CRO", "✅ 미팅완료"],
    ["⑥ CMC", "원료 품질관리, 안정성\n제형 스케일업", "병행 진행", "1~2억 (추정)", "미정", "⚠ 미탐색"],
    ["⑦ IND 준비", "비임상 패키지 통합\n식약처 사전상담/제출", "3~6개월", "RA/CRO 비용", "미정", "⚠ 미탐색"],
]

col_widths = [Inches(1.3), Inches(2.8), Inches(1.0), Inches(1.5), Inches(1.3), Inches(1.0)]
tbl_shape, tbl = make_table(slide1, Inches(0.4), Inches(4.0), Inches(8.9),
                             len(data_rows)+1, len(headers), col_widths)

for j, h in enumerate(headers):
    style_cell(tbl.cell(0, j), h, size=8, bold=True, color=WHITE, fill=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(data_rows):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    for j, val in enumerate(row):
        c = ACCENT_RED if val in ["미정"] else DARK_GRAY
        style_cell(tbl.cell(i+1, j), val, size=7, color=c, fill=bg)

# ── 비용 요약 + 핵심 포인트 ──
# 비용 박스
cost_box = add_shape(slide1, Inches(9.6), Inches(1.0), Inches(3.4), Inches(2.3), RGBColor(0xF0, 0xF4, 0xFA), LIGHT_BLUE, Pt(1.5))
tf = set_text(cost_box, "💰 비임상 예상 총 비용", size=12, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", size=4)
add_paragraph(tf, "최소 (4주 패키지)", size=9, bold=True, color=MID_BLUE)
add_paragraph(tf, "~10억 원", size=18, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", size=4)
add_paragraph(tf, "최대 (13주 패키지)", size=9, bold=True, color=ACCENT_ORANGE)
add_paragraph(tf, "~15억 원", size=18, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

# 핵심 포인트 박스
key_box = add_shape(slide1, Inches(9.6), Inches(3.5), Inches(3.4), Inches(3.6), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(1.5))
tf2 = set_text(key_box, "⚡ 핵심 포인트", size=12, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
points = [
    ("허가 트랙", "펩타이드 캐리어 = 신규 물질\n→ 신약 트랙 (Dt&CRO 확인)"),
    ("독성 범위", "캐리어 단독 + 복합체\n(캐리어+Finasteride) 모두 평가"),
    ("투자 유치 전", "최소 DRF(4주 반복) 독성\n데이터 확보 권장"),
    ("식약처 상담", "캐리어 분류 (부형제 vs 신물질)\n→ 시험 범위에 큰 영향"),
]
for title, desc in points:
    add_paragraph(tf2, "", size=3)
    add_paragraph(tf2, f"▸ {title}", size=9, bold=True, color=ACCENT_ORANGE)
    add_paragraph(tf2, f"  {desc}", size=7, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# 페이지 2: 임상 (Clinical) 파이프라인 + 협업 업체
# ═══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# ── 상단 바 ──
header_bar2 = add_rect(slide2, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
set_text(header_bar2, "임상 (Clinical) 파이프라인 & 협업 업체 현황", size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
header_bar2.text_frame.paragraphs[0].space_before = Pt(8)

# ── 임상 단계 테이블 ──
add_text_box(slide2, Inches(0.4), Inches(1.0), Inches(3), Inches(0.3),
             "▎ 임상 단계별 계획", size=13, bold=True, color=NAVY)

clin_headers = ["단계", "주요 내용", "기간", "예상 비용", "비고"]
clin_rows = [
    ["Phase 1", "안전성 + PK. 건강인/환자 20~80명\n외용제: 환자 직접 시험 가능 (1/2 통합)", "1~1.5년", "수억~10억+", "피부과 = 최저 비용군"],
    ["Phase 2", "AGA 환자 유효성/안전성 확인\n100~250명. 용량 탐색", "2~2.5년", "수십억 원", "성공률 ~29%\n(\"죽음의 계곡\")"],
    ["Phase 3", "대규모 확증 시험. 300~3,000명\n미녹시딜/경구 Finasteride 대조", "2.5~3년", "수백억 원", "다국가 임상 시\n비용 분산 가능"],
    ["NDA 허가", "품목허가 신청 → 식약처/FDA 심사", "1~2년", "RA 비용", "식약처 목표:\n295일→240일"],
]

clin_col_widths = [Inches(1.0), Inches(3.5), Inches(1.0), Inches(1.3), Inches(1.8)]
tbl2_shape, tbl2 = make_table(slide2, Inches(0.4), Inches(1.4), Inches(8.6),
                               len(clin_rows)+1, len(clin_headers), clin_col_widths)

for j, h in enumerate(clin_headers):
    style_cell(tbl2.cell(0, j), h, size=8, bold=True, color=WHITE, fill=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(clin_rows):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    for j, val in enumerate(row):
        style_cell(tbl2.cell(i+1, j), val, size=7, fill=bg)

# ── 외용제 이점 박스 ──
adv_box = add_shape(slide2, Inches(9.5), Inches(1.0), Inches(3.5), Inches(2.8), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(1.5))
tf_adv = set_text(adv_box, "🎯 국소 외용제 이점", size=12, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
advantages = [
    "Phase 1에서 환자 직접 투여 가능\n→ Phase 1/2 통합 설계로 기간 단축",
    "전신 노출 최소 → 일부 독성시험\n항목 면제 가능성 (식약처 협의)",
    "피부과 영역은 전체 적응증 중\n최저 비용군 (ASPE 보고서)",
    "국소 Finasteride FDA 미승인\n→ 최초 승인 시 시장 선점 기회",
]
for adv in advantages:
    add_paragraph(tf_adv, "", size=3)
    add_paragraph(tf_adv, f"✓ {adv}", size=7, color=DARK_GRAY)

# ── 협업 업체 현황 ──
add_text_box(slide2, Inches(0.4), Inches(4.05), Inches(4), Inches(0.3),
             "▎ 협업 업체 현황", size=13, bold=True, color=NAVY)

partner_headers = ["업체", "역할", "미팅일", "상태", "핵심 메모"]
partner_rows = [
    ["에피바이오텍", "비임상 효력 시험\n(탈모 전문 CRO)", "2025.12.23", "✅", "견적 받기로 함. IND 컨설팅 가능\n남성형 질환 모델·인간 모낭 시험 보유"],
    ["Dt&CRO", "GLP 독성 시험\n분석법 개발", "2025.11.26", "✅", "4주 패키지 ~8억 / 13주 ~10-12억\n13주 기준 견적 받기로 함"],
    ["P&K 피부임상연구센터", "인체적용시험\n(화장품/기능성)", "2025.12.23", "✅", "기능성 화장품 인증용 6개월~1년\n피부 흡수도 평가 가능"],
    ["Bio FD&C", "OEM/ODM\nINCI 원료등록", "2025.10~11", "✅", "INCI 등록 진행 (2026.02.09)\n화장품 제형 생산 담당"],
    ["OPIS (이탈리아)", "글로벌 CRO\n(Phase I~IV)", "2026.02", "✅", "30개국 네트워크\nEU 임상 협력 탐색"],
    ["인하대 신현태 교수", "임상 설계 자문", "2025.09", "✅", "임상 디자인 조언\nHCP 유통 전략"],
    ["BLT 특허법인", "IP + 사업화", "-", "✅", "국내외 특허 전략\n법인 설립 지원"],
    ["CMC 전문업체", "원료의약품 품질관리\n분석법, 안정성", "-", "⚠", "필수 확보 필요\nDt&CRO 강력 권고"],
    ["RA 전문가", "인허가 전략\n개발 플랜 수립", "-", "⚠", "필수 확보 필요\n\"산업의 눈\" (Dt&CRO 조언)"],
]

p_col_widths = [Inches(1.5), Inches(1.5), Inches(0.9), Inches(0.5), Inches(3.2)]
tbl3_shape, tbl3 = make_table(slide2, Inches(0.4), Inches(4.4), Inches(7.6),
                               len(partner_rows)+1, len(partner_headers), p_col_widths)

for j, h in enumerate(partner_headers):
    style_cell(tbl3.cell(0, j), h, size=8, bold=True, color=WHITE, fill=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(partner_rows):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    for j, val in enumerate(row):
        c = ACCENT_RED if val == "⚠" else (ACCENT_GREEN if val == "✅" else DARK_GRAY)
        style_cell(tbl3.cell(i+1, j), val, size=7, color=c, fill=bg,
                   alignment=PP_ALIGN.CENTER if j in [2,3] else PP_ALIGN.LEFT)

# ── 전체 로드맵 요약 (우측 하단) ──
road_box = add_shape(slide2, Inches(8.3), Inches(4.05), Inches(4.7), Inches(3.2), RGBColor(0xF5, 0xF0, 0xFA), MID_BLUE, Pt(1.5))
tf_road = set_text(road_box, "📋 전체 로드맵 (비임상 + 임상)", size=11, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

milestones = [
    ("2026 하반기", "법인 설립 + 제형 최적화 착수", "Seed"),
    ("2027", "효력 시험 + 분석법 개발 + CMC", "-"),
    ("2027~2029", "GLP 독성 시험 풀 패키지", "Pre-A"),
    ("2029", "IND 제출 → Phase 1 진입", "Series A"),
    ("2030~2032", "Phase 2 → Phase 3", "Series B+"),
    ("2033~", "NDA 허가 → 상용화", "-"),
]
for year, desc, inv_round in milestones:
    add_paragraph(tf_road, "", size=2)
    round_text = f"  [{inv_round}]" if inv_round != "-" else ""
    add_paragraph(tf_road, f"▸ {year}", size=8, bold=True, color=MID_BLUE)
    add_paragraph(tf_road, f"   {desc}{round_text}", size=7, color=DARK_GRAY)


# ── 저장 ──
output_dir = os.path.dirname(os.path.abspath(__file__))
timestamp = datetime.now().strftime("%Y%m%d_%H%M")
output_path = os.path.join(output_dir, f"치료제_비임상_임상_파이프라인_{timestamp}_클로드.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
